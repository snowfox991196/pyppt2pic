import glob
import os
import sys
import shutil

import ppt2gif
from pptx import Presentation, util

# 使用传入参数作为路径
# ppt_file = sys.argv[1]
# print(ppt_file)

# 这里输入ppt文件路径，使用双斜杠（\\）
ppt_file = "D:\\PYProject\\pyppt2pic\\test.pptx"

# 拆分路径，文件名，扩展名
filepath, tempfilename = os.path.split(ppt_file)
filename, extension = os.path.splitext(tempfilename)
print(filepath)  # 路径
print(filename)  # 文件名
print(extension)  # 扩展名

ppt_obj = ppt2gif.PPT(ppt_file)
ppt_obj.convert2png()

path_file_number = glob.glob(filepath + '\\' + filename + '\\*.png')

print(path_file_number)
print(len(path_file_number))

# 实例化一个ppt演示文稿对象
prs = Presentation()
# 调整页面大小
prs.slide_width = util.Cm(32)
prs.slide_height = util.Cm(18)

for i in range(len(path_file_number)):
    # 实例化空白模板
    blank_slide_layout = prs.slide_layouts[6]
    # 向文件中添加空白页面
    slide = prs.slides.add_slide(blank_slide_layout)

    img_path = filepath + '\\' + filename + '\\幻灯片' + str(i+1) + '.PNG'  # 图片路径
    # img_path = path_file_number[i]  # 图片路径由数组传入
    # 设置图片大小
    left = util.Cm(0)
    top = util.Cm(0)
    width = util.Cm(32)
    height = util.Cm(18)
    # 插入图片
    pic = prs.slides[i].shapes.add_picture(img_path, left, top, width, height)

# 保存为文件
prs.save(filepath + '\\' + filename + '-pic.pptx')
print("文件已保存至" + filepath + filename + '-pic.pptx')

# 删除临时文件夹
try:
    shutil.rmtree(filepath + '\\' + filename)
    print("已删除目录" + filepath + '\\' + filename)
except OSError as e:
    print("error %s - %s" % (e.filename, e.strerror))
